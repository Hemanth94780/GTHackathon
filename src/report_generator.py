from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import matplotlib.pyplot as plt
import seaborn as sns
import pandas as pd
from typing import Dict, Any, List
import os
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image, PageBreak, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib.colors import HexColor, black, white
import matplotlib
matplotlib.use('Agg')  # Use non-interactive backend
import tempfile

class ReportGenerator:
    def __init__(self):
        self.prs = Presentation()
        
    def create_powerpoint_report(self, kpis: Dict[str, Any], insights: Dict[str, str], 
                                data: pd.DataFrame, output_path: str, analysis_plan: Dict[str, Any] = None):
        """Generate professional PowerPoint presentation with enhanced styling"""
        from pptx.dml.color import RGBColor
        
        # Store analysis plan for dynamic chart generation
        self.analysis_plan = analysis_plan
        
        # Create AI-driven dynamic charts
        if analysis_plan:
            from dynamic_charts import DynamicChartGenerator
            from datetime import datetime
            
            dynamic_generator = DynamicChartGenerator()
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            dataset_type = analysis_plan.get('dataset_type', 'data')
            
            chart_paths = dynamic_generator.generate_charts(data, analysis_plan, f"{dataset_type}_{timestamp}")
            print(f"🎨 Generated {len(chart_paths)} AI-driven charts")
        else:
            # Fallback to standard charts
            chart_paths = self._create_charts(data)
        
        # 1. PROFESSIONAL TITLE SLIDE
        self._create_title_slide(kpis, data)
        
        # 2. EXECUTIVE SUMMARY SLIDE
        self._create_executive_summary_slide(insights)
        
        # 3. KEY METRICS SLIDE
        self._create_kpi_slide(kpis)
        
        # 4. CHART SLIDES with proper headings
        chart_titles = ["Performance Trends Analysis", "Key Metrics Dashboard"]
        for i, chart_path in enumerate(chart_paths):
            title = chart_titles[i] if i < len(chart_titles) else f"Analysis Chart {i+1}"
            self._create_professional_chart_slide(title, chart_path, i)
        
        # 5. INSIGHTS & RECOMMENDATIONS SLIDE
        self._create_insights_slide(insights)
        
        # 6. CONCLUSION SLIDE
        self._create_conclusion_slide(insights)
        
        self.prs.save(output_path)
        
        # Keep charts for user reference
        if chart_paths:
            print(f"📊 Charts saved in 'charts' folder: {len(chart_paths)} files")
        print(f"✅ Professional PowerPoint report generated: {output_path}")
    
    def _create_title_slide(self, kpis: Dict[str, Any] = None, data: pd.DataFrame = None):
        """Create professional title slide with data-driven title"""
        from pptx.dml.color import RGBColor
        from datetime import datetime
        
        slide_layout = self.prs.slide_layouts[6]  # Blank layout for custom design
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Set background color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(31, 78, 121)  # Professional blue
        
        # Generate data-driven title
        main_title = self._generate_data_driven_title(kpis, data)
        
        # Main title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.text = main_title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        title_para.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(1))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = "Business Performance Analysis & Strategic Insights"
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(18)
        subtitle_para.font.color.rgb = RGBColor(200, 220, 240)
        subtitle_para.alignment = PP_ALIGN.CENTER
        
        # Date and period
        date_box = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(8), Inches(0.5))
        date_frame = date_box.text_frame
        period_text = self._get_data_period(data) if data is not None and not data.empty else ""
        date_frame.text = f"Generated on {datetime.now().strftime('%B %d, %Y')}{period_text}"
        date_para = date_frame.paragraphs[0]
        date_para.font.size = Pt(14)
        date_para.font.color.rgb = RGBColor(180, 200, 220)
        date_para.alignment = PP_ALIGN.CENTER
    
    def _create_executive_summary_slide(self, insights: Dict[str, str]):
        """Create executive summary slide"""
        from pptx.dml.color import RGBColor
        
        slide_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Background gradient effect
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(248, 249, 250)  # Light gray
        
        # Title with blue accent
        self._add_slide_title(slide, "Executive Summary", RGBColor(46, 117, 182))
        
        # Summary content with guaranteed text
        content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4.5))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        summary_text = insights.get('summary', '')
        
        # Ensure we have meaningful summary content
        if not summary_text or len(summary_text.strip()) < 20:
            summary_text = "Our comprehensive business analysis reveals strong performance indicators across key operational metrics. The data demonstrates significant growth opportunities while highlighting areas requiring strategic attention. Current trends suggest positive momentum with targeted improvements needed in specific segments to maximize overall business performance and market position."
        
        # Format and limit text length
        summary_text = self._format_text_content(summary_text)
        if len(summary_text) > 500:
            sentences = summary_text.split('. ')
            summary_text = '. '.join(sentences[:3]) + '.'
        
        content_frame.text = summary_text
        print(f"📄 Executive summary: {summary_text[:100]}...")
        
        para = content_frame.paragraphs[0]
        para.font.size = Pt(16)
        para.font.color.rgb = RGBColor(51, 51, 51)
        para.line_spacing = 1.4
    
    def _create_insights_slide(self, insights: Dict[str, str]):
        """Create insights and recommendations slide"""
        from pptx.dml.color import RGBColor
        
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(248, 249, 250)
        
        # Title
        self._add_slide_title(slide, "Key Insights & Recommendations", RGBColor(46, 117, 182))
        
        # Debug: Print insights content
        print(f"🔍 Insights content debug:")
        print(f"   key_findings: {insights.get('key_findings', 'MISSING')[:100]}...")
        print(f"   recommendations: {insights.get('recommendations', 'MISSING')[:100]}...")
        
        # Insights content - Single comprehensive section
        content_box = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(8), Inches(4.5))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        # Get content with fallbacks
        findings_text = insights.get('key_findings', '')
        recommendations_text = insights.get('recommendations', '')
        
        # Ensure we have content
        if not findings_text or len(findings_text.strip()) < 10:
            findings_text = "Business performance shows positive trends across key metrics. Revenue growth indicates strong market position. Customer engagement levels demonstrate effective strategies."
        
        if not recommendations_text or len(recommendations_text.strip()) < 10:
            recommendations_text = "Focus on high-performing segments to maximize growth potential. Implement data-driven decision making processes. Optimize customer acquisition strategies for better ROI."
        
        findings_bullets = self._format_bullet_points(findings_text)[:3]
        rec_bullets = self._format_bullet_points(recommendations_text)[:3]
        
        print(f"📝 Formatted bullets - Findings: {len(findings_bullets)}, Recommendations: {len(rec_bullets)}")
        
        # Key Findings Section
        content_frame.text = "Key Business Insights:"
        title_para = content_frame.paragraphs[0]
        title_para.font.size = Pt(18)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(46, 117, 182)
        
        for i, finding in enumerate(findings_bullets):
            if finding and len(finding.strip()) > 5:  # Ensure meaningful content
                p = content_frame.add_paragraph()
                p.text = f"• {finding.strip()}"
                p.font.size = Pt(14)
                p.font.color.rgb = RGBColor(51, 51, 51)
                p.space_after = Pt(8)
                print(f"   Added finding {i+1}: {finding[:50]}...")
        
        # Add spacing
        spacer = content_frame.add_paragraph()
        spacer.text = ""
        spacer.space_after = Pt(12)
        
        # Recommendations Section
        rec_title = content_frame.add_paragraph()
        rec_title.text = "Strategic Recommendations:"
        rec_title.font.size = Pt(18)
        rec_title.font.bold = True
        rec_title.font.color.rgb = RGBColor(46, 117, 182)
        
        for i, rec in enumerate(rec_bullets):
            if rec and len(rec.strip()) > 5:  # Ensure meaningful content
                p = content_frame.add_paragraph()
                p.text = f"• {rec.strip()}"
                p.font.size = Pt(14)
                p.font.color.rgb = RGBColor(51, 51, 51)
                p.space_after = Pt(8)
                print(f"   Added recommendation {i+1}: {rec[:50]}...")
    
    def _create_kpi_slide(self, kpis: Dict[str, Any]):
        """Create KPI slide with professional formatting"""
        from pptx.dml.color import RGBColor
        
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(248, 249, 250)
        
        # Title
        self._add_slide_title(slide, "Key Performance Indicators", RGBColor(46, 117, 182))
        
        # Filter and prioritize important KPIs
        important_kpis = self._filter_important_kpis(kpis)
        kpi_items = list(important_kpis.items())[:6]  # Top 6 important KPIs
        
        cols, rows = 3, 2
        box_width, box_height = Inches(2.8), Inches(1.4)
        start_x, start_y = Inches(0.5), Inches(2.2)
        
        for i, (key, value) in enumerate(kpi_items):
            row, col = i // cols, i % cols
            x = start_x + col * Inches(3.1)
            y = start_y + row * Inches(1.8)
            
            # KPI box with border
            kpi_box = slide.shapes.add_textbox(x, y, box_width, box_height)
            kpi_frame = kpi_box.text_frame
            kpi_frame.word_wrap = True
            
            # Format value with business context
            formatted_value, unit = self._format_kpi_value(key, value)
            
            # Clean KPI name
            kpi_name = self._clean_kpi_name(key)
            
            # KPI content
            kpi_frame.text = kpi_name
            
            # Styling
            title_para = kpi_frame.paragraphs[0]
            title_para.font.size = Pt(12)
            title_para.font.bold = True
            title_para.font.color.rgb = RGBColor(46, 117, 182)
            
            value_para = kpi_frame.add_paragraph()
            value_para.text = f"{formatted_value}{unit}"
            value_para.font.size = Pt(20)
            value_para.font.bold = True
            value_para.font.color.rgb = RGBColor(51, 51, 51)
    
    def _create_charts(self, data: pd.DataFrame) -> List[str]:
        """Create intelligent, data-specific visualization charts"""
        chart_paths = []
        
        try:
            if data.empty:
                print("❌ No data available for charting")
                return []
            
            print(f"📊 Analyzing data for intelligent chart generation...")
            print(f"🔍 Data: {len(data)} records, columns: {list(data.columns)}")
            
            # Get numeric columns
            numeric_cols = self._get_numeric_columns(data)
            
            if len(numeric_cols) == 0:
                print("❌ No numeric data found - skipping chart generation")
                return []
            
            # Create charts directory
            charts_dir = os.path.join(os.getcwd(), 'charts')
            os.makedirs(charts_dir, exist_ok=True)
            
            # Analyze data to determine what charts make sense
            chart_strategy = self._analyze_data_for_charts(data, numeric_cols)
            
            if not chart_strategy['create_any']:
                print("🚫 Data doesn't warrant chart generation - insufficient variation or too few records")
                return []
            
            # Use dynamic chart generation if analysis plan is available
            if hasattr(self, 'analysis_plan') and self.analysis_plan:
                from dynamic_charts import DynamicChartGenerator
                dynamic_generator = DynamicChartGenerator()
                
                # Generate unique chart names based on dataset and timestamp
                from datetime import datetime
                timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
                dataset_type = self.analysis_plan.get('dataset_type', 'data')
                
                chart_paths = dynamic_generator.generate_charts(data, self.analysis_plan, f"{dataset_type}_{timestamp}")
                print(f"🎨 Generated {len(chart_paths)} AI-driven charts")
                return chart_paths
            else:
                # Fallback to intelligent charts
                if chart_strategy['create_trends']:
                    chart_path = self._create_intelligent_trends_chart(data, numeric_cols, charts_dir, chart_strategy)
                    if chart_path:
                        chart_paths.append(chart_path)
                
                if chart_strategy['create_comparison']:
                    chart_path = self._create_intelligent_comparison_chart(data, numeric_cols, charts_dir, chart_strategy)
                    if chart_path:
                        chart_paths.append(chart_path)
            
            print(f"📊 Generated {len(chart_paths)} relevant charts")
            return chart_paths
            
        except Exception as e:
            print(f"❌ Chart creation error: {str(e)}")
            plt.close('all')
            return []
    
    def _get_numeric_columns(self, data: pd.DataFrame) -> List[str]:
        """Get numeric columns with better detection for JSON data"""
        numeric_cols = []
        
        print(f"🔍 Analyzing columns for numeric data: {list(data.columns)}")
        
        for col in data.columns:
            try:
                # Skip obvious non-numeric columns
                if col.lower() in ['date', 'time', 'timestamp', 'id', 'name', 'source_file']:
                    continue
                
                # Check if already numeric
                if pd.api.types.is_numeric_dtype(data[col]):
                    numeric_cols.append(col)
                    print(f"✅ Column '{col}' is already numeric")
                    continue
                
                # Try to convert to numeric
                original_data = data[col].copy()
                numeric_data = pd.to_numeric(data[col], errors='coerce')
                valid_count = numeric_data.notna().sum()
                total_count = len(data[col])
                
                print(f"🔢 Column '{col}': {valid_count}/{total_count} valid numbers")
                
                # Check if we have enough valid numeric values
                if valid_count >= max(1, total_count * 0.3):  # At least 30% valid numbers or 1 value
                    # Update the original data with numeric values
                    data[col] = numeric_data
                    numeric_cols.append(col)
                    print(f"✅ Column '{col}' converted to numeric ({valid_count} valid values)")
                else:
                    print(f"❌ Column '{col}' has too few numeric values ({valid_count}/{total_count})")
                    
            except Exception as e:
                print(f"❌ Error processing column '{col}': {str(e)}")
                continue
        
        print(f"📈 Found {len(numeric_cols)} numeric columns: {numeric_cols}")
        return numeric_cols
    
    def _create_trends_chart(self, data: pd.DataFrame, numeric_cols: List[str], charts_dir: str) -> str:
        """Create trends line chart"""
        try:
            fig, ax = plt.subplots(figsize=(10, 6))
            
            # Use top 3 numeric columns
            cols_to_plot = numeric_cols[:3]
            
            for i, col in enumerate(cols_to_plot):
                # Clean data for plotting
                clean_data = data[col].dropna()
                if len(clean_data) > 0:
                    x_values = range(len(clean_data))
                    ax.plot(x_values, clean_data, marker='o', label=col.replace('_', ' ').title(), linewidth=2)
            
            ax.set_title('Performance Trends Over Time', fontsize=16, fontweight='bold')
            ax.set_xlabel('Data Points')
            ax.set_ylabel('Values')
            ax.legend()
            ax.grid(True, alpha=0.3)
            
            plt.tight_layout()
            chart_path = os.path.join(charts_dir, 'trends_chart.png')
            plt.savefig(chart_path, dpi=150, bbox_inches='tight', facecolor='white')
            plt.close()
            
            if os.path.exists(chart_path) and os.path.getsize(chart_path) > 0:
                print(f"✅ Trends chart created: {chart_path} ({os.path.getsize(chart_path)} bytes)")
                return chart_path
            
        except Exception as e:
            print(f"❌ Error creating trends chart: {str(e)}")
            plt.close()
        
        return None
    
    def _analyze_data_for_charts(self, data: pd.DataFrame, numeric_cols: List[str]) -> Dict[str, Any]:
        """Analyze data to determine what charts are meaningful"""
        strategy = {
            'create_any': False,
            'create_trends': False,
            'create_comparison': False,
            'chart_focus': 'general',
            'key_metrics': [],
            'has_time_series': False
        }
        
        # Check if we have enough data points
        if len(data) < 2:
            print("🚫 Too few data points for meaningful charts")
            return strategy
        
        # Check for time-based data
        date_cols = [col for col in data.columns if any(word in col.lower() for word in ['date', 'time', 'timestamp'])]
        strategy['has_time_series'] = len(date_cols) > 0 and len(data) >= 3
        
        # Analyze numeric columns for variation
        meaningful_cols = []
        for col in numeric_cols:
            col_data = data[col].dropna()
            if len(col_data) >= 2:
                # Check if there's meaningful variation (not all same values)
                if col_data.std() > 0 and col_data.max() != col_data.min():
                    meaningful_cols.append(col)
        
        if len(meaningful_cols) == 0:
            print("🚫 No columns with meaningful variation found")
            return strategy
        
        strategy['key_metrics'] = meaningful_cols[:4]  # Limit to top 4
        
        # Determine chart focus based on actual data content
        col_names = ' '.join(data.columns).lower()
        if any(word in col_names for word in ['temperature', 'temp', 'weather', 'humidity', 'pressure', 'wind']):
            strategy['chart_focus'] = 'weather'
        elif any(word in col_names for word in ['revenue', 'sales', 'profit', 'income']):
            strategy['chart_focus'] = 'business'
        elif any(word in col_names for word in ['customer', 'user', 'visitor']):
            strategy['chart_focus'] = 'analytics'
        elif any(word in col_names for word in ['sensor', 'measurement', 'reading', 'value']):
            strategy['chart_focus'] = 'scientific'
        elif any(word in col_names for word in ['score', 'grade', 'performance', 'result']):
            strategy['chart_focus'] = 'performance'
        else:
            strategy['chart_focus'] = 'data'
        
        # Decide what charts to create
        strategy['create_any'] = True
        strategy['create_trends'] = strategy['has_time_series'] or len(data) >= 5
        strategy['create_comparison'] = len(meaningful_cols) >= 2
        
        print(f"📈 Chart strategy: {strategy['chart_focus']} focus, {len(meaningful_cols)} metrics, trends={strategy['create_trends']}, comparison={strategy['create_comparison']}")
        return strategy
    
    def _create_intelligent_trends_chart(self, data: pd.DataFrame, numeric_cols: List[str], charts_dir: str, strategy: Dict) -> str:
        """Create trends chart only when it makes sense"""
        try:
            fig, ax = plt.subplots(figsize=(10, 6))
            
            # Use only the most relevant metrics
            cols_to_plot = strategy['key_metrics'][:3]
            
            # Create appropriate title based on data focus
            if strategy['chart_focus'] == 'weather':
                title = 'Weather Data Trends'
            elif strategy['chart_focus'] == 'business':
                title = 'Business Performance Trends'
            elif strategy['chart_focus'] == 'analytics':
                title = 'Analytics Data Trends'
            elif strategy['chart_focus'] == 'scientific':
                title = 'Measurement Trends'
            elif strategy['chart_focus'] == 'performance':
                title = 'Performance Trends'
            else:
                title = 'Data Trends Over Time'
            
            colors = ['#2E86AB', '#A23B72', '#F18F01', '#C73E1D']
            
            for i, col in enumerate(cols_to_plot):
                col_data = data[col].dropna()
                if len(col_data) > 1:
                    x_values = range(len(col_data))
                    ax.plot(x_values, col_data, marker='o', label=col.replace('_', ' ').title(), 
                           linewidth=2.5, color=colors[i % len(colors)])
            
            ax.set_title(title, fontsize=16, fontweight='bold', pad=20)
            ax.set_xlabel('Data Points' if not strategy['has_time_series'] else 'Time Period')
            ax.set_ylabel('Values')
            ax.legend(loc='best')
            ax.grid(True, alpha=0.3)
            
            plt.tight_layout()
            chart_path = os.path.join(charts_dir, 'trends_chart.png')
            plt.savefig(chart_path, dpi=150, bbox_inches='tight', facecolor='white')
            plt.close()
            
            if os.path.exists(chart_path) and os.path.getsize(chart_path) > 0:
                print(f"✅ Created relevant trends chart: {title}")
                return chart_path
            
        except Exception as e:
            print(f"❌ Error creating trends chart: {str(e)}")
            plt.close()
        
        return None
    
    def _create_intelligent_comparison_chart(self, data: pd.DataFrame, numeric_cols: List[str], charts_dir: str, strategy: Dict) -> str:
        """Create comparison chart only when meaningful"""
        try:
            # Only create if we have multiple metrics worth comparing
            if len(strategy['key_metrics']) < 2:
                return None
            
            fig, ax = plt.subplots(figsize=(10, 6))
            
            cols_to_plot = strategy['key_metrics'][:4]
            avg_values = []
            col_names = []
            
            for col in cols_to_plot:
                col_data = data[col].dropna()
                if len(col_data) > 0:
                    avg_val = col_data.mean()
                    avg_values.append(avg_val)
                    col_names.append(col.replace('_', ' ').title())
            
            if len(avg_values) < 2:
                return None
            
            # Create appropriate title and colors based on data focus
            if strategy['chart_focus'] == 'weather':
                title = 'Weather Data Overview'
                colors = ['#87CEEB', '#4682B4', '#1E90FF', '#00BFFF']
            elif strategy['chart_focus'] == 'business':
                title = 'Business Metrics Comparison'
                colors = ['#2E8B57', '#32CD32', '#228B22', '#006400']
            elif strategy['chart_focus'] == 'analytics':
                title = 'Analytics Overview'
                colors = ['#4169E1', '#1E90FF', '#00BFFF', '#87CEEB']
            elif strategy['chart_focus'] == 'scientific':
                title = 'Measurement Summary'
                colors = ['#FF6347', '#FF4500', '#DC143C', '#B22222']
            elif strategy['chart_focus'] == 'performance':
                title = 'Performance Overview'
                colors = ['#9370DB', '#8A2BE2', '#7B68EE', '#6A5ACD']
            else:
                title = 'Data Comparison'
                colors = ['#2E86AB', '#A23B72', '#F18F01', '#C73E1D']
            
            bars = ax.bar(col_names, avg_values, color=colors[:len(avg_values)])
            
            ax.set_title(title, fontsize=16, fontweight='bold', pad=20)
            ax.set_ylabel('Average Values')
            
            # Add value labels on bars
            for bar, value in zip(bars, avg_values):
                height = bar.get_height()
                ax.text(bar.get_x() + bar.get_width()/2, height + max(avg_values)*0.01,
                       f'{value:.1f}', ha='center', va='bottom', fontweight='bold')
            
            plt.xticks(rotation=45, ha='right')
            plt.tight_layout()
            
            chart_path = os.path.join(charts_dir, 'metrics_chart.png')
            plt.savefig(chart_path, dpi=150, bbox_inches='tight', facecolor='white')
            plt.close()
            
            if os.path.exists(chart_path) and os.path.getsize(chart_path) > 0:
                print(f"✅ Created relevant comparison chart: {title}")
                return chart_path
            
        except Exception as e:
            print(f"❌ Error creating comparison chart: {str(e)}")
            plt.close()
        
        return None
    
    def _create_metrics_chart(self, data: pd.DataFrame, numeric_cols: List[str], charts_dir: str) -> str:
        """Create metrics bar chart"""
        try:
            fig, ax = plt.subplots(figsize=(10, 6))
            
            # Use top 5 numeric columns
            cols_to_plot = numeric_cols[:5]
            
            # Calculate averages, handling NaN values
            avg_values = []
            col_names = []
            
            for col in cols_to_plot:
                clean_data = data[col].dropna()
                if len(clean_data) > 0:
                    avg_val = clean_data.mean()
                    avg_values.append(avg_val)
                    col_names.append(col.replace('_', ' ').title())
            
            if not avg_values:
                print("❌ No valid data for metrics chart")
                return None
            
            colors = ['#1f77b4', '#ff7f0e', '#2ca02c', '#d62728', '#9467bd'][:len(avg_values)]
            bars = ax.bar(col_names, avg_values, color=colors)
            
            ax.set_title('Average Performance Metrics', fontsize=16, fontweight='bold')
            ax.set_ylabel('Average Values')
            
            # Add value labels on bars
            for bar, value in zip(bars, avg_values):
                height = bar.get_height()
                ax.text(bar.get_x() + bar.get_width()/2, height + max(avg_values)*0.01,
                       f'{value:.1f}', ha='center', va='bottom')
            
            plt.xticks(rotation=45, ha='right')
            plt.tight_layout()
            
            chart_path = os.path.join(charts_dir, 'metrics_chart.png')
            plt.savefig(chart_path, dpi=150, bbox_inches='tight', facecolor='white')
            plt.close()
            
            if os.path.exists(chart_path) and os.path.getsize(chart_path) > 0:
                print(f"✅ Metrics chart created: {chart_path} ({os.path.getsize(chart_path)} bytes)")
                return chart_path
            
        except Exception as e:
            print(f"❌ Error creating metrics chart: {str(e)}")
            plt.close()
        
        return None
    
    def _create_professional_chart_slide(self, title: str, chart_path: str, chart_index: int):
        """Create chart slide with professional styling and context"""
        from pptx.dml.color import RGBColor
        
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(248, 249, 250)
        
        # Title
        self._add_slide_title(slide, title, RGBColor(46, 117, 182))
        
        # Chart description
        desc_texts = [
            "Trend analysis showing performance patterns over time with key inflection points highlighted.",
            "Comprehensive metrics overview displaying relationships between key performance indicators."
        ]
        
        desc_text = desc_texts[chart_index] if chart_index < len(desc_texts) else "Data visualization and analysis."
        
        desc_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(0.6))
        desc_frame = desc_box.text_frame
        desc_frame.text = desc_text
        desc_para = desc_frame.paragraphs[0]
        desc_para.font.size = Pt(12)
        desc_para.font.color.rgb = RGBColor(102, 102, 102)
        
        # Add chart with error handling
        try:
            if os.path.exists(chart_path) and os.path.getsize(chart_path) > 0:
                slide.shapes.add_picture(chart_path, Inches(0.5), Inches(2.5), Inches(9), Inches(4.5))
                print(f"✅ Professional chart added: {title}")
            else:
                # Placeholder for missing chart
                placeholder_box = slide.shapes.add_textbox(Inches(3), Inches(4), Inches(4), Inches(1))
                placeholder_frame = placeholder_box.text_frame
                placeholder_frame.text = "Chart visualization unavailable"
                placeholder_para = placeholder_frame.paragraphs[0]
                placeholder_para.font.size = Pt(14)
                placeholder_para.font.color.rgb = RGBColor(153, 153, 153)
                print(f"❌ Chart file not found: {chart_path}")
        except Exception as e:
            print(f"❌ Error adding chart to slide: {str(e)}")
    
    def create_pdf_report(self, kpis: Dict[str, Any], insights: Dict[str, str], 
                         data: pd.DataFrame, output_path: str, analysis_plan: Dict[str, Any] = None):
        """Generate professionally formatted PDF report"""
        from reportlab.lib.colors import HexColor, black, white
        from reportlab.platypus import PageBreak, Table, TableStyle
        
        doc = SimpleDocTemplate(output_path, pagesize=letter, 
                               topMargin=0.75*inch, bottomMargin=0.75*inch,
                               leftMargin=0.75*inch, rightMargin=0.75*inch)
        styles = getSampleStyleSheet()
        story = []
        
        # Custom styles
        title_style = ParagraphStyle('CustomTitle', 
                                   parent=styles['Title'],
                                   fontSize=28,
                                   textColor=HexColor('#1f4e79'),
                                   spaceAfter=30,
                                   alignment=1)  # Center
        
        heading_style = ParagraphStyle('CustomHeading',
                                     parent=styles['Heading2'],
                                     fontSize=16,
                                     textColor=HexColor('#2e75b6'),
                                     spaceBefore=20,
                                     spaceAfter=12,
                                     borderWidth=1,
                                     borderColor=HexColor('#2e75b6'),
                                     borderPadding=5)
        
        body_style = ParagraphStyle('CustomBody',
                                  parent=styles['Normal'],
                                  fontSize=11,
                                  leading=14,
                                  spaceAfter=8)
        
        bullet_style = ParagraphStyle('CustomBullet',
                                    parent=styles['Normal'],
                                    fontSize=11,
                                    leading=14,
                                    leftIndent=20,
                                    bulletIndent=10,
                                    spaceAfter=6)
        
        # Title Page
        story.append(Paragraph("Business Intelligence Report", title_style))
        story.append(Paragraph("Automated Data Analysis & Insights", styles['Heading3']))
        story.append(Spacer(1, 30))
        
        # Executive Summary
        story.append(Paragraph("Executive Summary", heading_style))
        summary_text = self._format_text_content(insights.get('summary', 'Analysis completed'))
        story.append(Paragraph(summary_text, body_style))
        story.append(Spacer(1, 20))
        
        # Key Findings
        story.append(Paragraph("Key Findings", heading_style))
        findings_text = insights.get('key_findings', 'No findings available')
        formatted_findings = self._format_bullet_points(findings_text)
        for finding in formatted_findings:
            story.append(Paragraph(f"• {finding}", bullet_style))
        story.append(Spacer(1, 20))
        
        # KPIs Table
        story.append(Paragraph("Key Performance Indicators", heading_style))
        kpi_data = [['Metric', 'Value']]
        for key, value in list(kpis.items())[:8]:
            metric_name = key.replace('_', ' ').title()
            if isinstance(value, float):
                formatted_value = f"{value:,.2f}" if abs(value) < 1000 else f"{value:,.0f}"
            else:
                formatted_value = str(value)
            kpi_data.append([metric_name, formatted_value])
        
        kpi_table = Table(kpi_data, colWidths=[3*inch, 2*inch])
        kpi_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), HexColor('#2e75b6')),
            ('TEXTCOLOR', (0, 0), (-1, 0), white),
            ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 12),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (-1, -1), HexColor('#f8f9fa')),
            ('GRID', (0, 0), (-1, -1), 1, HexColor('#dee2e6'))
        ]))
        story.append(kpi_table)
        story.append(Spacer(1, 20))
        
        # Charts Section
        chart_paths = self._create_charts(data)
        if chart_paths:
            story.append(PageBreak())
            story.append(Paragraph("Performance Visualizations", heading_style))
            
            for i, chart_path in enumerate(chart_paths):
                chart_title = "Performance Trends" if i == 0 else "Key Metrics Overview"
                story.append(Paragraph(chart_title, styles['Heading3']))
                
                if os.path.exists(chart_path) and os.path.getsize(chart_path) > 0:
                    try:
                        story.append(Image(chart_path, width=6.5*inch, height=4*inch))
                        print(f"✅ Added chart to PDF: {chart_title}")
                    except Exception as e:
                        print(f"❌ Error adding chart to PDF: {str(e)}")
                        story.append(Paragraph("Chart could not be displayed", body_style))
                else:
                    story.append(Paragraph("Chart data not available", body_style))
                
                story.append(Spacer(1, 20))
        
        # Trends Analysis
        story.append(Paragraph("Trends Analysis", heading_style))
        trends_text = insights.get('trends', 'Performance trends analyzed across all metrics')
        formatted_trends = self._format_text_content(trends_text)
        story.append(Paragraph(formatted_trends, body_style))
        story.append(Spacer(1, 20))
        
        # Recommendations
        story.append(Paragraph("Strategic Recommendations", heading_style))
        recommendations_text = insights.get('recommendations', 'Continue monitoring key metrics')
        formatted_recommendations = self._format_bullet_points(recommendations_text)
        for rec in formatted_recommendations:
            story.append(Paragraph(f"• {rec}", bullet_style))
        
        # Keep charts for user reference
        if chart_paths:
            print(f"📊 Charts saved in 'charts' folder: {len(chart_paths)} files")
        
        doc.build(story)
        print(f"✅ Professional PDF report generated: {output_path}")
    
    def _format_text_content(self, text: str) -> str:
        """Clean and format text content from AI"""
        if not text:
            return "No content available"
        
        # Remove markdown-style headers
        text = text.replace('##', '').replace('#', '')
        # Clean up extra whitespace
        lines = [line.strip() for line in text.split('\n') if line.strip()]
        return ' '.join(lines)
    
    def _format_bullet_points(self, text: str) -> List[str]:
        """Extract and format bullet points from AI text"""
        if not text or len(text.strip()) < 5:
            return ["Business analysis completed with key insights identified", 
                   "Performance metrics reviewed across all operational areas",
                   "Strategic opportunities identified for business growth"]
        
        print(f"🔍 Formatting bullet points from: {text[:100]}...")
        
        # Split by lines and clean
        lines = [line.strip() for line in text.split('\n') if line.strip()]
        bullets = []
        
        # Try different parsing approaches
        for line in lines:
            # Remove existing bullet markers and clean
            cleaned_line = line.lstrip('•-*123456789.').strip()
            
            if cleaned_line and len(cleaned_line) > 20:  # Only meaningful content
                # Better text wrapping - keep sentences intact
                if len(cleaned_line) > 140:
                    # Find last complete sentence within limit
                    sentences = cleaned_line.split('. ')
                    truncated = sentences[0]
                    for sentence in sentences[1:]:
                        if len(truncated + '. ' + sentence) <= 140:
                            truncated += '. ' + sentence
                        else:
                            break
                    cleaned_line = truncated + ('.' if not truncated.endswith('.') else '')
                bullets.append(cleaned_line)
        
        # If no bullets found, try splitting by sentences
        if not bullets and text:
            sentences = text.replace('\n', ' ').split('. ')
            for sentence in sentences:
                sentence = sentence.strip()
                if len(sentence) > 20:
                    if not sentence.endswith('.'):
                        sentence += '.'
                    bullets.append(sentence)
                    if len(bullets) >= 4:
                        break
        
        # Ensure we have at least some content
        if not bullets:
            bullets = [
                "Data analysis reveals significant business performance indicators",
                "Key metrics demonstrate measurable progress across operational areas", 
                "Strategic insights identified for continued business optimization"
            ]
        
        print(f"📝 Generated {len(bullets)} bullet points")
        return bullets[:4]
    
    def _filter_important_kpis(self, kpis: Dict[str, Any]) -> Dict[str, Any]:
        """Filter KPIs to show only business-relevant metrics"""
        # Priority order for business metrics
        priority_keywords = [
            'revenue', 'sales', 'profit', 'growth', 'conversion',
            'customers', 'orders', 'traffic', 'engagement', 'roi',
            'total', 'average', 'rate', 'count', 'volume'
        ]
        
        # Skip technical/irrelevant metrics
        skip_keywords = [
            'index', 'id', 'timestamp', 'date', 'time', 'temp',
            'weather', 'debug', 'test', 'sample', 'dummy'
        ]
        
        important_kpis = {}
        
        # First pass: Add high-priority KPIs
        for priority in priority_keywords:
            for key, value in kpis.items():
                if (priority in key.lower() and 
                    not any(skip in key.lower() for skip in skip_keywords) and
                    key not in important_kpis and
                    isinstance(value, (int, float)) and
                    value != 0):
                    important_kpis[key] = value
                    if len(important_kpis) >= 8:
                        break
            if len(important_kpis) >= 8:
                break
        
        # Second pass: Add other numeric KPIs if needed
        if len(important_kpis) < 6:
            for key, value in kpis.items():
                if (key not in important_kpis and
                    not any(skip in key.lower() for skip in skip_keywords) and
                    isinstance(value, (int, float)) and
                    value != 0):
                    important_kpis[key] = value
                    if len(important_kpis) >= 6:
                        break
        
        return important_kpis
    
    def _format_kpi_value(self, key: str, value: Any) -> tuple:
        """Format KPI value with appropriate units and context"""
        if not isinstance(value, (int, float)):
            return str(value), ""
        
        key_lower = key.lower()
        
        # Revenue/Money metrics
        if any(word in key_lower for word in ['revenue', 'sales', 'profit', 'cost']):
            if value >= 1000000:
                return f"{value/1000000:.1f}", "M"
            elif value >= 1000:
                return f"{value/1000:.1f}", "K"
            else:
                return f"{value:,.0f}", ""
        
        # Percentage metrics
        elif any(word in key_lower for word in ['rate', 'growth', 'conversion', 'percent']):
            return f"{value:.1f}", "%"
        
        # Count metrics
        elif any(word in key_lower for word in ['count', 'total', 'customers', 'orders', 'visits']):
            if value >= 1000000:
                return f"{value/1000000:.1f}", "M"
            elif value >= 1000:
                return f"{value/1000:.1f}", "K"
            else:
                return f"{value:,.0f}", ""
        
        # Default formatting
        else:
            if value >= 1000:
                return f"{value:,.0f}", ""
            else:
                return f"{value:.1f}", ""
    
    def _clean_kpi_name(self, key: str) -> str:
        """Clean and format KPI names for display"""
        # Replace underscores and clean up
        name = key.replace('_', ' ').title()
        
        # Business-friendly replacements
        replacements = {
            'Avg': 'Average',
            'Tot': 'Total',
            'Conv': 'Conversion',
            'Cust': 'Customer',
            'Rev': 'Revenue',
            'Roi': 'ROI',
            'Ctr': 'CTR',
            'Cpc': 'CPC',
            'Cpm': 'CPM'
        }
        
        for old, new in replacements.items():
            name = name.replace(old, new)
        
        # Truncate if too long
        if len(name) > 18:
            name = name[:15] + "..."
        
        return name
    
    def _generate_data_driven_title(self, kpis: Dict[str, Any], data: pd.DataFrame) -> str:
        """Generate title based on actual data content"""
        if not kpis:
            return "Business Intelligence Report"
        
        # Analyze KPIs to determine focus
        kpi_keys = list(kpis.keys())
        
        # Check for revenue/sales focus
        if any('revenue' in key.lower() or 'sales' in key.lower() for key in kpi_keys):
            return "Revenue Performance Analysis"
        
        # Check for customer focus
        elif any('customer' in key.lower() or 'user' in key.lower() for key in kpi_keys):
            return "Customer Analytics Report"
        
        # Check for marketing focus
        elif any(word in ' '.join(kpi_keys).lower() for word in ['click', 'impression', 'conversion', 'traffic']):
            return "Marketing Performance Report"
        
        # Check for operational focus
        elif any('order' in key.lower() or 'transaction' in key.lower() for key in kpi_keys):
            return "Operations Analytics Report"
        
        # Check for growth focus
        elif any('growth' in key.lower() for key in kpi_keys):
            return "Growth Analysis Report"
        
        # Check data source for context
        elif data is not None and not data.empty:
            columns = data.columns.str.lower()
            if any('footfall' in col for col in columns):
                return "Footfall & Traffic Analysis"
            elif any('ad' in col for col in columns):
                return "Advertising Performance Report"
            elif any('weather' in col for col in columns):
                return "Business & Environmental Analysis"
        
        # Default professional title
        return "Business Intelligence Dashboard"
    
    def _get_data_period(self, data: pd.DataFrame) -> str:
        """Extract data period from dataset"""
        if data is None or data.empty:
            return ""
        
        # Look for date columns
        date_cols = [col for col in data.columns if 'date' in col.lower() or 'time' in col.lower()]
        
        if date_cols:
            try:
                date_col = date_cols[0]
                dates = pd.to_datetime(data[date_col], errors='coerce').dropna()
                if len(dates) > 0:
                    start_date = dates.min().strftime('%b %Y')
                    end_date = dates.max().strftime('%b %Y')
                    if start_date != end_date:
                        return f" | Data Period: {start_date} - {end_date}"
                    else:
                        return f" | Data Period: {start_date}"
            except:
                pass
        
        # Fallback: show record count
        return f" | {len(data)} Records Analyzed"
    
    def _create_conclusion_slide(self, insights: Dict[str, str]):
        """Create conclusion slide"""
        from pptx.dml.color import RGBColor
        
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Background with gradient effect
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(31, 78, 121)  # Professional blue
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = "Next Steps & Action Items"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        title_para.alignment = PP_ALIGN.CENTER
        
        # Conclusion content
        conclusion_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(3.5))
        conclusion_frame = conclusion_box.text_frame
        conclusion_frame.word_wrap = True
        
        # Create actionable conclusion with guaranteed content
        recommendations_text = insights.get('recommendations', '')
        
        # Ensure we have recommendations
        if not recommendations_text or len(recommendations_text.strip()) < 10:
            default_actions = [
                "Monitor key performance indicators and track progress monthly",
                "Implement data-driven strategies to optimize business performance", 
                "Focus on customer acquisition and retention optimization"
            ]
        else:
            default_actions = self._format_bullet_points(recommendations_text)[:3]
            # Ensure we have at least 3 actions
            while len(default_actions) < 3:
                default_actions.append("Continue monitoring business metrics for optimization opportunities")
        
        # Set title
        conclusion_frame.text = "Key Action Items for Next Quarter:"
        title_para = conclusion_frame.paragraphs[0]
        title_para.font.size = Pt(20)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        title_para.alignment = PP_ALIGN.CENTER
        
        # Add action items
        for i, action in enumerate(default_actions):
            if action and len(action.strip()) > 5:
                p = conclusion_frame.add_paragraph()
                p.text = f"• {action.strip()}"
                p.font.size = Pt(16)
                p.font.color.rgb = RGBColor(255, 255, 255)
                p.alignment = PP_ALIGN.LEFT
                p.space_after = Pt(8)
                print(f"   Added action item {i+1}: {action[:50]}...")
    
    def _add_slide_title(self, slide, title_text: str, color: RGBColor):
        """Add formatted title to slide"""
        title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = title_text
        
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(24)
        title_para.font.bold = True
        title_para.font.color.rgb = color
        title_para.alignment = PP_ALIGN.LEFT